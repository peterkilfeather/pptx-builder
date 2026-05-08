"""
Reusable diagram generators for python-pptx.
Each function adds native, editable PowerPoint shapes to a slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def extract_template_colors(pptx_path):
    """Extract theme colors from a template. Returns a dict of color roles to RGBColor values."""
    prs = Presentation(pptx_path)
    theme = prs.slide_masters[0].slide_layouts[0].slide_master.element
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    colors = {}
    color_map = {
        "dk1": "dark1",
        "dk2": "dark2",
        "lt1": "light1",
        "lt2": "light2",
        "accent1": "accent1",
        "accent2": "accent2",
        "accent3": "accent3",
        "accent4": "accent4",
        "accent5": "accent5",
        "accent6": "accent6",
    }

    for elem in theme.iter():
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        if tag in color_map:
            srgb = elem.find("a:srgbClr", ns)
            if srgb is not None:
                hex_val = srgb.get("val")
                colors[color_map[tag]] = RGBColor.from_string(hex_val)
            sys_clr = elem.find("a:sysClr", ns)
            if sys_clr is not None:
                hex_val = sys_clr.get("lastClr")
                if hex_val:
                    colors[color_map[tag]] = RGBColor.from_string(hex_val)

    return colors


def _add_shape(slide, shape_type, left, top, width, height, text="",
               fill_color=None, border_color=None, border_width=Pt(1),
               font_size=Pt(11), font_color=None, bold=False, alignment=PP_ALIGN.CENTER):
    """Add a shape with text and styling."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        p.font.size = font_size
        p.font.bold = bold
        if font_color:
            p.font.color.rgb = font_color

    return shape


def _add_connector(slide, start_x, start_y, end_x, end_y,
                   color=None, width=Pt(1.5)):
    """Add a straight connector line between two points."""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, start_x, start_y, end_x, end_y
    )
    if color:
        connector.line.color.rgb = color
    connector.line.width = width
    return connector


def _add_arrow(slide, start_x, start_y, end_x, end_y,
               color=None, width=Pt(1.5)):
    """Add a connector line with an arrowhead."""
    connector = _add_connector(slide, start_x, start_y, end_x, end_y, color, width)
    end_el = connector.line._ln
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    from lxml import etree
    tail_end = etree.SubElement(end_el, f"{{{ns}}}tailEnd")
    tail_end.set("type", "triangle")
    tail_end.set("w", "med")
    tail_end.set("len", "med")
    return connector


def create_flowchart(slide, steps, colors=None, top_margin=Inches(1.2)):
    """
    Create a top-to-bottom flowchart with process boxes and decision diamonds.

    Args:
        slide: pptx slide object
        steps: list of dicts, each with:
            - "text": label text
            - "type": "process" | "decision" | "start" | "end"
        colors: dict with keys "fill", "border", "text", "arrow" as RGBColor
        top_margin: space reserved at top for slide title
    """
    if colors is None:
        colors = {
            "fill": RGBColor(0x06, 0x5A, 0x82),
            "border": RGBColor(0x04, 0x3B, 0x56),
            "text": RGBColor(0xFF, 0xFF, 0xFF),
            "arrow": RGBColor(0x36, 0x45, 0x4F),
            "decision_fill": RGBColor(0xF9, 0x61, 0x67),
        }

    shape_map = {
        "process": MSO_SHAPE.RECTANGLE,
        "decision": MSO_SHAPE.DIAMOND,
        "start": MSO_SHAPE.ROUNDED_RECTANGLE,
        "end": MSO_SHAPE.ROUNDED_RECTANGLE,
    }

    box_w = Inches(2.2)
    box_h = Inches(0.6)
    diamond_w = Inches(2.0)
    diamond_h = Inches(1.1)
    gap = Inches(0.35)
    n = len(steps)
    bottom_margin = Inches(0.5)

    total_height = sum(
        (diamond_h if s["type"] == "decision" else box_h) for s in steps
    ) + gap * (n - 1)
    available = Inches(7.5) - top_margin - bottom_margin
    start_y = top_margin + (available - total_height) // 2
    start_y = max(start_y, top_margin)
    center_x = Inches(5.0)

    prev_bottom = None
    prev_center_x = None

    for step in steps:
        is_diamond = step["type"] == "decision"
        w = diamond_w if is_diamond else box_w
        h = diamond_h if is_diamond else box_h
        left = center_x - w // 2
        top = start_y

        fill = colors.get("decision_fill", colors["fill"]) if is_diamond else colors["fill"]

        _add_shape(
            slide, shape_map[step["type"]], left, top, w, h,
            text=step["text"],
            fill_color=fill,
            border_color=colors.get("border"),
            font_size=Pt(11),
            font_color=colors.get("text"),
            bold=True,
        )

        if prev_bottom is not None:
            _add_arrow(
                slide,
                prev_center_x, prev_bottom,
                center_x, top,
                color=colors.get("arrow"),
            )

        prev_bottom = top + h
        prev_center_x = center_x
        start_y = top + h + gap


def create_process_flow(slide, steps, colors=None):
    """
    Create a horizontal process flow with arrows between steps.

    Args:
        slide: pptx slide object
        steps: list of strings (step labels)
        colors: dict with "fill", "border", "text", "arrow" as RGBColor
    """
    if colors is None:
        colors = {
            "fill": RGBColor(0x06, 0x5A, 0x82),
            "border": RGBColor(0x04, 0x3B, 0x56),
            "text": RGBColor(0xFF, 0xFF, 0xFF),
            "arrow": RGBColor(0x36, 0x45, 0x4F),
        }

    n = len(steps)
    margin = Inches(0.7)
    arrow_gap = Inches(0.5)
    total_width = Inches(10.0) - 2 * margin
    box_w = (total_width - arrow_gap * (n - 1)) // n
    box_h = Inches(0.9)
    top = Inches(3.0)

    for i, text in enumerate(steps):
        left = margin + i * (box_w + arrow_gap)

        _add_shape(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h,
            text=text,
            fill_color=colors["fill"],
            border_color=colors.get("border"),
            font_size=Pt(12),
            font_color=colors.get("text"),
            bold=True,
        )

        if i < n - 1:
            arrow_start_x = left + box_w
            arrow_end_x = left + box_w + arrow_gap
            arrow_y = top + box_h // 2
            _add_arrow(
                slide, arrow_start_x, arrow_y, arrow_end_x, arrow_y,
                color=colors.get("arrow"),
            )


def create_timeline(slide, events, colors=None):
    """
    Create a horizontal timeline with labeled milestones.

    Args:
        slide: pptx slide object
        events: list of dicts with "date" and "label"
        colors: dict with "line", "dot", "text" as RGBColor
    """
    if colors is None:
        colors = {
            "line": RGBColor(0x36, 0x45, 0x4F),
            "dot": RGBColor(0x06, 0x5A, 0x82),
            "text": RGBColor(0x21, 0x21, 0x21),
            "date": RGBColor(0x06, 0x5A, 0x82),
        }

    n = len(events)
    margin = Inches(1.3)
    line_y = Inches(3.8)
    line_width = Inches(7.4)
    dot_size = Inches(0.25)

    _add_connector(
        slide,
        margin, line_y,
        margin + line_width, line_y,
        color=colors["line"],
        width=Pt(2),
    )

    spacing = line_width // max(n - 1, 1)

    for i, event in enumerate(events):
        cx = margin + spacing * i if n > 1 else margin + line_width // 2
        dot_left = cx - dot_size // 2
        dot_top = line_y - dot_size // 2

        _add_shape(
            slide, MSO_SHAPE.OVAL, dot_left, dot_top, dot_size, dot_size,
            fill_color=colors["dot"],
        )

        alternating = i % 2 == 0
        label_top = line_y - Inches(1.2) if alternating else line_y + Inches(0.4)
        date_top = line_y - Inches(0.6) if alternating else line_y + Inches(0.9)

        label_w = Inches(1.5)
        label_left = cx - label_w // 2

        _add_shape(
            slide, MSO_SHAPE.RECTANGLE, label_left, label_top, label_w, Inches(0.4),
            text=event["label"],
            font_size=Pt(10),
            font_color=colors["text"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        _add_shape(
            slide, MSO_SHAPE.RECTANGLE, label_left, date_top, label_w, Inches(0.3),
            text=event.get("date", ""),
            font_size=Pt(9),
            font_color=colors["date"],
            alignment=PP_ALIGN.CENTER,
        )


def create_comparison(slide, columns, colors=None):
    """
    Create side-by-side comparison columns.

    Args:
        slide: pptx slide object
        columns: list of dicts with "title" and "points" (list of strings)
        colors: dict with "header_fill", "header_text", "body_fill", "body_text" as RGBColor
    """
    if colors is None:
        colors = {
            "header_fill": RGBColor(0x06, 0x5A, 0x82),
            "header_text": RGBColor(0xFF, 0xFF, 0xFF),
            "body_fill": RGBColor(0xF2, 0xF2, 0xF2),
            "body_text": RGBColor(0x21, 0x21, 0x21),
        }

    n = len(columns)
    margin = Inches(0.7)
    col_gap = Inches(0.4)
    total_width = Inches(10.0) - 2 * margin
    col_w = (total_width - col_gap * (n - 1)) // n
    header_h = Inches(0.6)
    body_h = Inches(3.5)
    top = Inches(1.8)

    for i, col in enumerate(columns):
        left = margin + i * (col_w + col_gap)

        _add_shape(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, header_h,
            text=col["title"],
            fill_color=colors["header_fill"],
            font_size=Pt(14),
            font_color=colors["header_text"],
            bold=True,
        )

        body = _add_shape(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top + header_h + Inches(0.1),
            col_w, body_h,
            fill_color=colors["body_fill"],
            border_color=RGBColor(0xDD, 0xDD, 0xDD),
        )

        tf = body.text_frame
        tf.word_wrap = True
        from pptx.enum.text import MSO_ANCHOR
        tf.paragraphs[0].space_before = Pt(8)
        body.text_frame.auto_size = None
        body.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        for j, point in enumerate(col.get("points", [])):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"  {point}"
            p.font.size = Pt(12)
            p.font.color.rgb = colors["body_text"]
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)


def create_hierarchy(slide, tree, colors=None, _x=None, _y=None, _width=None):
    """
    Create a top-down hierarchy / org chart.

    Args:
        slide: pptx slide object
        tree: dict with "text" and optional "children" (list of tree dicts)
        colors: dict with "fill", "border", "text", "line" as RGBColor
    """
    if colors is None:
        colors = {
            "fill": RGBColor(0x06, 0x5A, 0x82),
            "border": RGBColor(0x04, 0x3B, 0x56),
            "text": RGBColor(0xFF, 0xFF, 0xFF),
            "line": RGBColor(0x36, 0x45, 0x4F),
        }

    if _x is None:
        _x = Inches(0.5)
        _y = Inches(1.5)
        _width = Inches(9.0)

    box_w = min(Inches(2.0), _width - Inches(0.3))
    box_h = Inches(0.6)
    level_gap = Inches(1.0)

    center_x = _x + _width // 2
    box_left = max(center_x - box_w // 2, _x + Inches(0.1))

    _add_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE, box_left, _y, box_w, box_h,
        text=tree["text"],
        fill_color=colors["fill"],
        border_color=colors.get("border"),
        font_size=Pt(11),
        font_color=colors["text"],
        bold=True,
    )

    children = tree.get("children", [])
    if not children:
        return

    child_width = _width // len(children)
    child_y = _y + box_h + level_gap

    for i, child in enumerate(children):
        child_x = _x + i * child_width
        child_center_x = child_x + child_width // 2

        _add_arrow(
            slide,
            center_x, _y + box_h,
            child_center_x, child_y,
            color=colors["line"],
        )

        create_hierarchy(
            slide, child, colors,
            _x=child_x, _y=child_y, _width=child_width,
        )
